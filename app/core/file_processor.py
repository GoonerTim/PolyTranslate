"""File processor for reading various document formats."""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import TYPE_CHECKING, Any

import chardet
import pandas as pd
import pptx
import pypdf
from bs4 import BeautifulSoup
from docx import Document
from markdown import markdown

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)


class FileProcessor:
    """Handles reading and processing of various file formats."""

    SUPPORTED_EXTENSIONS = {
        "txt",
        "pdf",
        "docx",
        "doc",
        "pptx",
        "xlsx",
        "xls",
        "csv",
        "html",
        "htm",
        "md",
        "markdown",
        "rpy",
        "srt",
        "ass",
        "ssa",
    }

    @staticmethod
    def detect_encoding(file_content: bytes) -> str:
        try:
            result = chardet.detect(file_content)
            detected = result.get("encoding")
            confidence = result.get("confidence", 0)

            if not detected or confidence < 0.7:
                try:
                    file_content.decode("utf-8")
                    return "utf-8"
                except UnicodeDecodeError:
                    pass

                if file_content.startswith(b"\xef\xbb\xbf"):
                    return "utf-8-sig"

                for enc in ["cp1251", "cp1252", "windows-1251", "windows-1252"]:
                    try:
                        file_content.decode(enc)
                        return enc
                    except (UnicodeDecodeError, LookupError):
                        continue

            return detected or "utf-8"
        except Exception:
            return "utf-8"

    @staticmethod
    def read_txt(file_content: bytes) -> str:
        encodings = [
            FileProcessor.detect_encoding(file_content),
            "utf-8",
            "utf-8-sig",
            "cp1251",
            "windows-1251",
            "cp1252",
            "windows-1252",
            "latin-1",
        ]

        for encoding in encodings:
            try:
                decoded = file_content.decode(encoding)
                if not all(ord(char) > 127 and ord(char) < 160 for char in decoded[:100]):
                    return decoded
            except (UnicodeDecodeError, LookupError, AttributeError):
                continue

        return file_content.decode("utf-8", errors="replace")

    @staticmethod
    def read_pdf(file_content: bytes) -> str:
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = pypdf.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text
        except Exception as e:
            raise ValueError(f"PDF reading error: {e}") from e

    @staticmethod
    def read_docx(file_content: bytes) -> str:
        try:
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"DOCX reading error: {e}") from e

    @staticmethod
    def read_pptx(file_content: bytes) -> str:
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = pptx.Presentation(pptx_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"PPTX reading error: {e}") from e

    @staticmethod
    def read_xlsx(file_content: bytes) -> str:
        try:
            xlsx_file = io.BytesIO(file_content)
            df = pd.read_excel(xlsx_file, sheet_name=None)
            text = ""
            for sheet_name, sheet_data in df.items():
                text += f"--- {sheet_name} ---\n"
                for _, row in sheet_data.iterrows():
                    row_text = " | ".join(str(cell) for cell in row if pd.notna(cell))
                    text += row_text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"XLSX reading error: {e}") from e

    @staticmethod
    def read_csv(file_content: bytes) -> str:
        try:
            encoding = FileProcessor.detect_encoding(file_content)
            csv_file = io.BytesIO(file_content)
            df = pd.read_csv(csv_file, encoding=encoding)
            text = ""
            for _, row in df.iterrows():
                row_text = " | ".join(str(cell) for cell in row if pd.notna(cell))
                text += row_text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"CSV reading error: {e}") from e

    @staticmethod
    def read_html(file_content: bytes) -> str:
        try:
            encoding = FileProcessor.detect_encoding(file_content)
            html_content = file_content.decode(encoding)
            soup = BeautifulSoup(html_content, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = " ".join(chunk for chunk in chunks if chunk)
            return text
        except Exception as e:
            raise ValueError(f"HTML reading error: {e}") from e

    @staticmethod
    def read_md(file_content: bytes) -> str:
        try:
            encoding = FileProcessor.detect_encoding(file_content)
            md_content = file_content.decode(encoding)
            html_content = markdown(md_content)
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text()
            return text
        except Exception as e:
            raise ValueError(f"Markdown reading error: {e}") from e

    # --- Delegated to SubtitleProcessor / RenpyProcessor ---
    # These static methods delegate to extracted classes for backwards compatibility.

    @staticmethod
    def read_rpy(
        file_content: bytes,
        translate_dialogue: bool = True,
        translate_strings: bool = True,
        preserve_code: bool = True,
    ) -> str:
        from app.core.renpy_processor import RenpyProcessor

        return RenpyProcessor.read_rpy(
            file_content, translate_dialogue, translate_strings, preserve_code
        )

    @staticmethod
    def reconstruct_rpy(
        original_content: str,
        translations: dict[str, str],
        translate_dialogue: bool = True,
        translate_strings: bool = True,
    ) -> str:
        from app.core.renpy_processor import RenpyProcessor

        return RenpyProcessor.reconstruct_rpy(
            original_content, translations, translate_dialogue, translate_strings
        )

    @staticmethod
    def split_rpy_by_scenes(content: str) -> list[tuple[str, str]]:
        from app.core.renpy_processor import RenpyProcessor

        return RenpyProcessor.split_rpy_by_scenes(content)

    @staticmethod
    def read_srt(file_content: bytes) -> str:
        from app.core.subtitle_processor import SubtitleProcessor

        return SubtitleProcessor.read_srt(file_content)

    @staticmethod
    def reconstruct_srt(original_content: str, translations: dict[str, str]) -> str:
        from app.core.subtitle_processor import SubtitleProcessor

        return SubtitleProcessor.reconstruct_srt(original_content, translations)

    @staticmethod
    def read_ass(file_content: bytes) -> str:
        from app.core.subtitle_processor import SubtitleProcessor

        return SubtitleProcessor.read_ass(file_content)

    @staticmethod
    def reconstruct_ass(original_content: str, translations: dict[str, str]) -> str:
        from app.core.subtitle_processor import SubtitleProcessor

        return SubtitleProcessor.reconstruct_ass(original_content, translations)

    _PROCESSORS: dict[str, str] = {
        "txt": "read_txt",
        "pdf": "read_pdf",
        "docx": "read_docx",
        "doc": "read_docx",
        "pptx": "read_pptx",
        "xlsx": "read_xlsx",
        "xls": "read_xlsx",
        "csv": "read_csv",
        "html": "read_html",
        "htm": "read_html",
        "md": "read_md",
        "markdown": "read_md",
        "rpy": "read_rpy",
        "srt": "read_srt",
        "ass": "read_ass",
        "ssa": "read_ass",
    }

    @classmethod
    def process_file(cls, file_path: str | Path, **kwargs: Any) -> str:
        path = Path(file_path)
        extension = path.suffix.lower().lstrip(".")

        if extension not in cls.SUPPORTED_EXTENSIONS:
            logger.warning("Unsupported extension '%s', trying as plain text", extension)
            try:
                content = path.read_bytes()
                return cls.read_txt(content)
            except Exception as e:
                raise ValueError(f"Unsupported file format: {extension}") from e

        content = path.read_bytes()

        method_name = cls._PROCESSORS.get(extension)
        if method_name is None:
            return cls.read_txt(content)

        processor = getattr(cls, method_name)
        if extension == "rpy":
            return processor(content, **kwargs)
        return processor(content)

    @classmethod
    def process_bytes(cls, file_content: bytes, extension: str, **kwargs: Any) -> str:
        extension = extension.lower().lstrip(".")

        method_name = cls._PROCESSORS.get(extension, "read_txt")
        processor = getattr(cls, method_name)
        if extension == "rpy":
            return processor(file_content, **kwargs)
        return processor(file_content)
