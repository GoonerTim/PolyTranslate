"""Additional tests for file processor formats."""

from __future__ import annotations

import io
from pathlib import Path

import pandas as pd
import pptx
import pypdf
from docx import Document

from app.core.file_processor import FileProcessor


class TestFileProcessorPDF:
    """Tests for PDF file processing."""

    def test_read_pdf_simple(self) -> None:
        # Create a simple PDF in memory
        pdf_writer = pypdf.PdfWriter()
        pdf_writer.add_blank_page(width=200, height=200)

        # Add text to the page
        output = io.BytesIO()
        pdf_writer.write(output)
        pdf_content = output.getvalue()

        # Should not raise an error
        result = FileProcessor.read_pdf(pdf_content)
        assert isinstance(result, str)

    def test_read_pdf_empty(self) -> None:
        pdf_writer = pypdf.PdfWriter()
        pdf_writer.add_blank_page(width=200, height=200)

        output = io.BytesIO()
        pdf_writer.write(output)
        pdf_content = output.getvalue()

        result = FileProcessor.read_pdf(pdf_content)
        assert isinstance(result, str)

    def test_read_pdf_invalid(self) -> None:
        invalid_content = b"Not a PDF file"
        try:
            FileProcessor.read_pdf(invalid_content)
        except ValueError as e:
            assert "PDF reading error" in str(e)


class TestFileProcessorDOCX:
    """Tests for DOCX file processing."""

    def test_read_docx_simple(self) -> None:
        # Create a simple DOCX in memory
        doc = Document()
        doc.add_paragraph("Hello, world!")
        doc.add_paragraph("This is a test document.")

        output = io.BytesIO()
        doc.save(output)
        docx_content = output.getvalue()

        result = FileProcessor.read_docx(docx_content)
        assert "Hello, world!" in result
        assert "This is a test document." in result

    def test_read_docx_empty(self) -> None:
        doc = Document()
        output = io.BytesIO()
        doc.save(output)
        docx_content = output.getvalue()

        result = FileProcessor.read_docx(docx_content)
        assert isinstance(result, str)

    def test_read_docx_multiple_paragraphs(self) -> None:
        doc = Document()
        for i in range(5):
            doc.add_paragraph(f"Paragraph {i}")

        output = io.BytesIO()
        doc.save(output)
        docx_content = output.getvalue()

        result = FileProcessor.read_docx(docx_content)
        for i in range(5):
            assert f"Paragraph {i}" in result

    def test_read_docx_invalid(self) -> None:
        invalid_content = b"Not a DOCX file"
        try:
            FileProcessor.read_docx(invalid_content)
        except ValueError as e:
            assert "DOCX reading error" in str(e)


class TestFileProcessorPPTX:
    """Tests for PPTX file processing."""

    def test_read_pptx_simple(self) -> None:
        # Create a simple PPTX in memory
        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"

        output = io.BytesIO()
        presentation.save(output)
        pptx_content = output.getvalue()

        result = FileProcessor.read_pptx(pptx_content)
        assert "Test Presentation" in result

    def test_read_pptx_multiple_slides(self) -> None:
        presentation = pptx.Presentation()

        for i in range(3):
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            title = slide.shapes.title
            title.text = f"Slide {i}"

        output = io.BytesIO()
        presentation.save(output)
        pptx_content = output.getvalue()

        result = FileProcessor.read_pptx(pptx_content)
        for i in range(3):
            assert f"Slide {i}" in result

    def test_read_pptx_empty(self) -> None:
        presentation = pptx.Presentation()
        output = io.BytesIO()
        presentation.save(output)
        pptx_content = output.getvalue()

        result = FileProcessor.read_pptx(pptx_content)
        assert isinstance(result, str)

    def test_read_pptx_invalid(self) -> None:
        invalid_content = b"Not a PPTX file"
        try:
            FileProcessor.read_pptx(invalid_content)
        except ValueError as e:
            assert "PPTX reading error" in str(e)


class TestFileProcessorXLSX:
    """Tests for XLSX file processing."""

    def test_read_xlsx_simple(self) -> None:
        # Create a simple XLSX in memory
        df = pd.DataFrame({"A": [1, 2, 3], "B": ["x", "y", "z"]})

        output = io.BytesIO()
        df.to_excel(output, index=False, sheet_name="Sheet1", engine="openpyxl")
        xlsx_content = output.getvalue()

        result = FileProcessor.read_xlsx(xlsx_content)
        assert "1" in result or "x" in result

    def test_read_xlsx_multiple_sheets(self) -> None:
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            df1 = pd.DataFrame({"A": [1, 2]})
            df2 = pd.DataFrame({"B": [3, 4]})
            df1.to_excel(writer, sheet_name="Sheet1", index=False)
            df2.to_excel(writer, sheet_name="Sheet2", index=False)

        xlsx_content = output.getvalue()
        result = FileProcessor.read_xlsx(xlsx_content)
        assert "Sheet1" in result
        assert "Sheet2" in result

    def test_read_xlsx_empty(self) -> None:
        df = pd.DataFrame()
        output = io.BytesIO()
        df.to_excel(output, index=False, engine="openpyxl")
        xlsx_content = output.getvalue()

        result = FileProcessor.read_xlsx(xlsx_content)
        assert isinstance(result, str)

    def test_read_xlsx_invalid(self) -> None:
        invalid_content = b"Not an XLSX file"
        try:
            FileProcessor.read_xlsx(invalid_content)
        except ValueError as e:
            assert "XLSX reading error" in str(e)


class TestFileProcessorCSV:
    """Tests for CSV file processing."""

    def test_read_csv_simple(self) -> None:
        csv_content = b"Name,Age\nAlice,30\nBob,25"
        result = FileProcessor.read_csv(csv_content)
        assert "Alice" in result
        assert "Bob" in result
        assert "30" in result
        assert "25" in result

    def test_read_csv_with_encoding(self) -> None:
        csv_content = "Имя,Возраст\nАлиса,30\nБоб,25".encode()
        result = FileProcessor.read_csv(csv_content)
        assert "Алиса" in result or "30" in result

    def test_read_csv_empty(self) -> None:
        csv_content = b""
        try:
            result = FileProcessor.read_csv(csv_content)
            # May succeed with empty string or raise error
            assert isinstance(result, str)
        except ValueError:
            # Empty CSV may cause error
            pass

    def test_read_csv_invalid(self) -> None:
        invalid_content = b"\x00\x01\x02\x03"  # Binary garbage
        try:
            FileProcessor.read_csv(invalid_content)
        except ValueError as e:
            assert "CSV reading error" in str(e)


class TestFileProcessorIntegration:
    """Integration tests for file processor."""

    def test_process_file_auto_detect(self, temp_dir: Path) -> None:
        # Create a test file
        file_path = temp_dir / "test.txt"
        file_path.write_text("Hello, world!", encoding="utf-8")

        result = FileProcessor.process_file(file_path)
        assert result == "Hello, world!"

    def test_process_bytes_auto_detect(self) -> None:
        content = b"Test content"
        result = FileProcessor.process_bytes(content, "txt")
        assert result == "Test content"

    def test_process_bytes_rpy(self) -> None:
        content = b'e "Hello"\nmc "World"'
        result = FileProcessor.process_bytes(
            content, "rpy", translate_dialogue=True, translate_strings=False
        )
        assert "Hello" in result

    def test_process_file_unsupported_extension(self, temp_dir: Path) -> None:
        file_path = temp_dir / "test.unknown"
        file_path.write_text("Some content", encoding="utf-8")

        # Should fall back to text reading
        result = FileProcessor.process_file(file_path)
        assert result == "Some content"

    def test_supported_extensions_complete(self) -> None:
        expected = {
            "txt",
            "pdf",
            "docx",
            "doc",
            "pptx",
            "xlsx",
            "xls",
            "csv",
            "html",
            "htm",
            "md",
            "markdown",
            "rpy",
            "srt",
            "ass",
            "ssa",
        }
        assert expected == FileProcessor.SUPPORTED_EXTENSIONS


class TestFileProcessorSRT:
    """Tests for SRT subtitle file processing."""

    def test_read_srt_simple(self) -> None:
        srt_content = (
            b"1\n00:00:01,000 --> 00:00:04,000\nHello, how are you?\n\n"
            b"2\n00:00:05,000 --> 00:00:08,000\nI'm fine, thanks.\n"
        )
        result = FileProcessor.read_srt(srt_content)
        assert "SRT_1: Hello, how are you?" in result
        assert "SRT_2: I'm fine, thanks." in result

    def test_read_srt_multiline_text(self) -> None:
        srt_content = (
            b"1\n00:00:01,000 --> 00:00:04,000\nLine one\nLine two\n\n"
            b"2\n00:00:05,000 --> 00:00:08,000\nSingle line\n"
        )
        result = FileProcessor.read_srt(srt_content)
        assert "SRT_1: Line one\nLine two" in result
        assert "SRT_2: Single line" in result

    def test_read_srt_empty_blocks_skipped(self) -> None:
        srt_content = (
            b"1\n00:00:01,000 --> 00:00:04,000\n\n\n2\n00:00:05,000 --> 00:00:08,000\nHello\n"
        )
        result = FileProcessor.read_srt(srt_content)
        assert "SRT_2: Hello" in result

    def test_read_srt_invalid(self) -> None:
        srt_content = b"Not a subtitle file at all"
        result = FileProcessor.read_srt(srt_content)
        # Falls back to returning original content
        assert "Not a subtitle file at all" in result

    def test_reconstruct_srt(self) -> None:
        original = (
            "1\n00:00:01,000 --> 00:00:04,000\nHello\n\n2\n00:00:05,000 --> 00:00:08,000\nWorld\n"
        )
        translations = {
            "SRT_1: Hello": "Привет",
            "SRT_2: World": "Мир",
        }
        result = FileProcessor.reconstruct_srt(original, translations)
        assert "Привет" in result
        assert "Мир" in result
        assert "00:00:01,000 --> 00:00:04,000" in result
        assert "00:00:05,000 --> 00:00:08,000" in result

    def test_reconstruct_srt_preserves_untranslated(self) -> None:
        original = "1\n00:00:01,000 --> 00:00:04,000\nHello\n"
        result = FileProcessor.reconstruct_srt(original, {})
        assert "Hello" in result
        assert "00:00:01,000 --> 00:00:04,000" in result

    def test_process_file_srt(self, temp_dir: Path) -> None:
        file_path = temp_dir / "test.srt"
        file_path.write_text("1\n00:00:01,000 --> 00:00:04,000\nHello\n", encoding="utf-8")
        result = FileProcessor.process_file(file_path)
        assert "SRT_1: Hello" in result

    def test_process_bytes_srt(self) -> None:
        content = b"1\n00:00:01,000 --> 00:00:04,000\nHello\n"
        result = FileProcessor.process_bytes(content, "srt")
        assert "SRT_1: Hello" in result


class TestFileProcessorASS:
    """Tests for ASS/SSA subtitle file processing."""

    SAMPLE_ASS = (
        "[Script Info]\nTitle: Test\n\n"
        "[V4+ Styles]\nFormat: Name, Fontname, Fontsize\n"
        "Style: Default,Arial,20\n\n"
        "[Events]\n"
        "Format: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n"
        "Dialogue: 0,0:00:01.00,0:00:04.00,Default,,0,0,0,,Hello world\n"
        "Dialogue: 0,0:00:05.00,0:00:08.00,Default,,0,0,0,,How are you?\n"
    )

    def test_read_ass_simple(self) -> None:
        result = FileProcessor.read_ass(self.SAMPLE_ASS.encode())
        assert "ASS_1: Hello world" in result
        assert "ASS_2: How are you?" in result

    def test_read_ass_with_override_tags(self) -> None:
        ass_content = (
            "[Events]\n"
            "Format: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n"
            r"Dialogue: 0,0:00:01.00,0:00:04.00,Default,,0,0,0,,{\b1}Bold text{\b0}" + "\n"
        ).encode()
        result = FileProcessor.read_ass(ass_content)
        assert "ASS_1:" in result
        assert "Bold text" in result

    def test_read_ass_skips_comments(self) -> None:
        ass_content = (
            b"[Events]\n"
            b"Format: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n"
            b"Comment: 0,0:00:01.00,0:00:04.00,Default,,0,0,0,,This is a comment\n"
            b"Dialogue: 0,0:00:05.00,0:00:08.00,Default,,0,0,0,,Visible text\n"
        )
        result = FileProcessor.read_ass(ass_content)
        assert "This is a comment" not in result
        assert "ASS_1: Visible text" in result

    def test_read_ass_text_with_commas(self) -> None:
        ass_content = (
            b"[Events]\n"
            b"Format: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n"
            b"Dialogue: 0,0:00:01.00,0:00:04.00,Default,,0,0,0,,Hello, world, how are you?\n"
        )
        result = FileProcessor.read_ass(ass_content)
        assert "ASS_1: Hello, world, how are you?" in result

    def test_read_ass_empty_events(self) -> None:
        ass_content = b"[Script Info]\nTitle: Test\n\n[Events]\nFormat: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n"
        result = FileProcessor.read_ass(ass_content)
        # No dialogues, returns original
        assert "[Script Info]" in result

    def test_reconstruct_ass(self) -> None:
        translations = {
            "ASS_1: Hello world": "Привет мир",
            "ASS_2: How are you?": "Как дела?",
        }
        result = FileProcessor.reconstruct_ass(self.SAMPLE_ASS, translations)
        assert "Привет мир" in result
        assert "Как дела?" in result
        assert "[Script Info]" in result
        assert "[V4+ Styles]" in result

    def test_reconstruct_ass_preserves_untranslated(self) -> None:
        result = FileProcessor.reconstruct_ass(self.SAMPLE_ASS, {})
        assert "Hello world" in result
        assert "How are you?" in result

    def test_process_file_ass(self, temp_dir: Path) -> None:
        file_path = temp_dir / "test.ass"
        file_path.write_text(self.SAMPLE_ASS, encoding="utf-8")
        result = FileProcessor.process_file(file_path)
        assert "ASS_1: Hello world" in result

    def test_process_file_ssa(self, temp_dir: Path) -> None:
        file_path = temp_dir / "test.ssa"
        file_path.write_text(self.SAMPLE_ASS, encoding="utf-8")
        result = FileProcessor.process_file(file_path)
        assert "ASS_1: Hello world" in result

    def test_process_bytes_ass(self) -> None:
        result = FileProcessor.process_bytes(self.SAMPLE_ASS.encode(), "ass")
        assert "ASS_1: Hello world" in result
